import os
import json
import secrets
import smtplib
from email.message import EmailMessage
import re
import sqlite3
import uuid
import base64
import mimetypes
import csv
import json
import io

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from datetime import datetime, timedelta
from functools import wraps
from urllib.parse import quote

import requests
from huggingface_hub import InferenceClient
from flask import Flask, request, jsonify, render_template, session, redirect, url_for, g, send_from_directory
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from pypdf import PdfReader
from authlib.integrations.flask_client import OAuth
from google.oauth2 import id_token as google_id_token
from google.auth.transport import requests as google_auth_requests

try:
    import psycopg2
    import psycopg2.extras
    PSYCOPG2_AVAILABLE = True
except ImportError:
    PSYCOPG2_AVAILABLE = False

# =========================
# Novara — Flask backend
# =========================

DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")

DEEPINFRA_API_KEY = os.getenv("DEEPINFRA_API_KEY", "")
DEEPINFRA_MODEL = os.getenv(
    "DEEPINFRA_MODEL",
    "meta-llama/Meta-Llama-3.1-8B-Instruct"
)


def call_deepinfra(prompt_text, timeout=60):
    if not DEEPINFRA_API_KEY:
        return None, "DEEPINFRA_API_KEY not set on server."

    try:
        response = requests.post(
            "https://api.deepinfra.com/v1/openai/chat/completions",
            headers={
                "Authorization": f"Bearer {DEEPINFRA_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": DEEPINFRA_MODEL,
                "messages": [
                    {"role": "user", "content": prompt_text}
                ],
            },
            timeout=timeout,
        )

        if response.status_code != 200:
            return None, f"DeepInfra error {response.status_code}: {response.text[:500]}"

        data = response.json()
        choices = data.get("choices", [])

        if not choices:
            return None, "DeepInfra returned no choices."

        content = choices[0].get("message", {}).get("content")
        if not content:
            return None, "DeepInfra returned empty content."

        return content, None

    except Exception as e:
        return None, f"DeepInfra request failed: {e}"


def call_deepseek(prompt_text, timeout=60):
    if not DEEPSEEK_API_KEY:
        return None, "DEEPSEEK_API_KEY not set on server."

    try:
        response = requests.post(
            "https://api.deepseek.com/chat/completions",
            headers={
                "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": "deepseek-chat",
                "messages": [{"role": "user", "content": prompt_text}],
                "stream": False
            },
            timeout=timeout
        )
        data = response.json()
        if "choices" not in data:
            err = data.get("error", {})
            return None, err.get("message", str(data)) if isinstance(err, dict) else str(data)
        return data["choices"][0]["message"]["content"], None
    except Exception as e:
        return None, str(e)
ADMIN_SECRET = os.getenv("ADMIN_SECRET", "")
GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID", "")
GOOGLE_CLIENT_SECRET = os.getenv("GOOGLE_CLIENT_SECRET", "")
DATABASE_URL = os.getenv("DATABASE_URL", "")  # set this to use permanent PostgreSQL instead of local SQLite
# Optional: plug in a real video-generation provider later (see README "Video generation").
VIDEO_API_URL = os.getenv("VIDEO_API_URL", "")
VIDEO_API_KEY = os.getenv("VIDEO_API_KEY", "")
MODEL = "deepseek-v4-flash"
VISION_MODEL = "deepseek-v4-flash-vision-exp"
HF_TOKEN = os.getenv("HF_TOKEN", "")
HF_IMAGE_MODEL = os.getenv("HF_IMAGE_MODEL", "stabilityai/sdxl-turbo")
HF_VIDEO_MODEL = "Wan-AI/Wan2.2-TI2V-5B"

# =========================
# NOVARA MODEL ROUTER
# =========================
NOVARA_FAST_MODEL = os.getenv(
    "NOVARA_FAST_MODEL",
    "meta-llama/Llama-3.1-8B-Instruct:cheapest"
)

NOVARA_THINKING_MODEL = os.getenv(
    "NOVARA_THINKING_MODEL",
    "google/gemma-3-12b-it:cheapest"
)

# Omega keeps using the existing direct DeepSeek connection.
# During the configured peak window, every mode can fall back to Gemma.
NOVARA_PEAK_START = int(os.getenv("NOVARA_PEAK_START", "19"))
NOVARA_PEAK_END = int(os.getenv("NOVARA_PEAK_END", "23"))

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(BASE_DIR, "novara.db")
UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

app = Flask(__name__)
app.secret_key = os.getenv("FLASK_SECRET", "change-this-secret-in-production")
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024  # 100MB max upload

IMAGE_EXT = {"png", "jpg", "jpeg", "webp", "gif", "bmp", "heic"}
VIDEO_EXT = {"mp4", "mov", "webm", "mkv", "avi"}
DOC_EXT = {"pdf"}
DOCX_EXT = {"docx"}
PPTX_EXT = {"pptx"}
XLSX_EXT = {"xlsx"}
TEXT_EXT = {"txt", "csv", "json"}
PLUGIN_EXT = {"doc", "docx", "ppt", "pptx", "xls", "xlsx", "txt", "csv", "json", "zip", "apk"}

# ---- Google OAuth setup ----
oauth = OAuth(app)
if GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET:
    oauth.register(
        name="google",
        client_id=GOOGLE_CLIENT_ID,
        client_secret=GOOGLE_CLIENT_SECRET,
        server_metadata_url="https://accounts.google.com/.well-known/openid-configuration",
        client_kwargs={"scope": "openid email profile"},
    )


# =========================
# DATABASE
# =========================

class PGConnWrapper:
    """Makes a psycopg2 connection behave like sqlite3's connection.execute()
    shorthand, so the rest of the app's code (written for SQLite) doesn't
    need to change. Also translates '?' placeholders to '%s' automatically."""

    def __init__(self, conn):
        self._conn = conn

    def execute(self, query, params=()):
        cur = self._conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)
        cur.execute(query.replace("?", "%s"), params)
        return cur

    def commit(self):
        self._conn.commit()

    def close(self):
        self._conn.close()


def get_db():
    if "db" not in g:
        if DATABASE_URL and PSYCOPG2_AVAILABLE:
            raw_conn = psycopg2.connect(DATABASE_URL)
            g.db = PGConnWrapper(raw_conn)
        else:
            g.db = sqlite3.connect(DB_PATH)
            g.db.row_factory = sqlite3.Row
    return g.db


def using_postgres():
    return bool(DATABASE_URL and PSYCOPG2_AVAILABLE)


@app.teardown_appcontext
def close_db(exception):
    db = g.pop("db", None)
    if db is not None:
        db.close()


def _add_column_if_missing(conn, table, column, coltype):
    if using_postgres():
        # Postgres supports this directly — no need to check first.
        conn.execute(f"ALTER TABLE {table} ADD COLUMN IF NOT EXISTS {column} {coltype}")
    else:
        cols = [r[1] for r in conn.execute(f"PRAGMA table_info({table})").fetchall()]
        if column not in cols:
            conn.execute(f"ALTER TABLE {table} ADD COLUMN {column} {coltype}")


def init_db():
    if using_postgres():
        raw_conn = psycopg2.connect(DATABASE_URL)
        conn = PGConnWrapper(raw_conn)
    else:
        conn = sqlite3.connect(DB_PATH)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS users (
            id TEXT PRIMARY KEY,
            username TEXT UNIQUE,
            email TEXT UNIQUE,
            google_id TEXT UNIQUE,
            password_hash TEXT,
            terms_accepted INTEGER DEFAULT 0,
            created_at TEXT NOT NULL
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS conversations (
            id TEXT PRIMARY KEY,
            user_id TEXT NOT NULL,
            title TEXT NOT NULL,
            share_id TEXT,
            is_shared INTEGER DEFAULT 0,
            created_at TEXT NOT NULL
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS messages (
            id TEXT PRIMARY KEY,
            conversation_id TEXT NOT NULL,
            role TEXT NOT NULL,
            text TEXT NOT NULL,
            attachment_path TEXT,
            attachment_type TEXT,
            sources TEXT,
            feedback TEXT,
            created_at TEXT NOT NULL
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS memories (
            id TEXT PRIMARY KEY,
            user_id TEXT NOT NULL,
            content TEXT NOT NULL,
            source TEXT,
            created_at TEXT NOT NULL
        )
    """)
    # Backfill columns for databases created before these features existed.
    _add_column_if_missing(conn, "users", "terms_accepted", "INTEGER DEFAULT 0")
    _add_column_if_missing(conn, "users", "phone_number", "TEXT")

    # =========================
    # NOVARA SUBSCRIPTION SYSTEM
    # =========================
    _add_column_if_missing(conn, "users", "plan", "TEXT DEFAULT 'free'")
    _add_column_if_missing(conn, "users", "subscription_product_id", "TEXT")
    _add_column_if_missing(conn, "users", "subscription_expiry", "TEXT")
    _add_column_if_missing(conn, "users", "subscription_purchase_token", "TEXT")

    _add_column_if_missing(conn, "users", "omega_enabled", "INTEGER DEFAULT 0")
    _add_column_if_missing(conn, "users", "image_used", "INTEGER DEFAULT 0")
    _add_column_if_missing(conn, "users", "silent_video_used", "INTEGER DEFAULT 0")
    _add_column_if_missing(conn, "users", "audio_video_used", "INTEGER DEFAULT 0")
    _add_column_if_missing(conn, "users", "voice_used", "INTEGER DEFAULT 0")
    _add_column_if_missing(conn, "users", "rewarded_today", "INTEGER DEFAULT 0")
    _add_column_if_missing(conn, "users", "rewarded_window_start", "TEXT")
    _add_column_if_missing(conn, "users", "rewarded_image_credits", "INTEGER DEFAULT 0")
    _add_column_if_missing(conn, "users", "rewarded_audio_video_credits", "INTEGER DEFAULT 0")
    _add_column_if_missing(conn, "users", "usage_month", "TEXT")

    _add_column_if_missing(conn, "conversations", "share_id", "TEXT")
    _add_column_if_missing(conn, "conversations", "is_shared", "INTEGER DEFAULT 0")
    _add_column_if_missing(conn, "messages", "feedback", "TEXT")

    # Existing users stay on the Free plan.
    conn.execute("""
        UPDATE users
        SET plan = 'free'
        WHERE plan IS NULL OR TRIM(plan) = ''
    """)

    conn.commit()
    conn.close()


init_db()



# =========================
# GOOGLE PLAY BILLING
# =========================

NOVARA_GOOGLE_PRODUCTS = {
    "novara_plus": "plus",
    "novara_pro": "pro",
    "novara_v32": "v3.2",
}

def verify_google_play_subscription(package_name, purchase_token, product_id):
    """Verify a Google Play subscription server-side."""
    try:
        import os
        import json
        from google.oauth2 import service_account
        from googleapiclient.discovery import build

        raw = os.environ.get(
            "GOOGLE_PLAY_SERVICE_ACCOUNT_JSON", ""
        ).strip()

        if not raw:
            raise RuntimeError(
                "Google Play service account is not configured."
            )

        credentials_info = json.loads(raw)

        credentials = service_account.Credentials.from_service_account_info(
            credentials_info,
            scopes=[
                "https://www.googleapis.com/auth/androidpublisher"
            ],
        )

        service = build(
            "androidpublisher",
            "v3",
            credentials=credentials,
            cache_discovery=False,
        )

        result = (
            service.purchases()
            .subscriptionsv2()
            .get(
                packageName=package_name,
                token=purchase_token,
            )
            .execute()
        )

        line_items = result.get("lineItems", [])

        matched = next(
            (
                item for item in line_items
                if item.get("productId") == product_id
            ),
            None,
        )

        if not matched:
            return {
                "verified": False,
                "error": "Google Play product does not match.",
            }

        state = result.get("subscriptionState", "")

        if state not in {
            "SUBSCRIPTION_STATE_ACTIVE",
            "SUBSCRIPTION_STATE_IN_GRACE_PERIOD",
        }:
            return {
                "verified": False,
                "error": "Google Play subscription is not active.",
            }

        return {
            "verified": True,
            "plan": NOVARA_GOOGLE_PRODUCTS[product_id],
            "product_id": product_id,
            "expiry": matched.get("expiryTime"),
            "order_id": result.get("latestOrderId"),
        }

    except Exception as exc:
        app.logger.exception("Google Play verification failed")
        return {
            "verified": False,
            "error": str(exc),
        }

# =========================
# AUTH HELPERS
# =========================

def login_required(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login_page"))
        return f(*args, **kwargs)
    return wrapper


def terms_required(f):
    """Blocks access until the user has ticked the Terms & Conditions checkbox."""
    @wraps(f)
    def wrapper(*args, **kwargs):
        db = get_db()
        user = db.execute("SELECT terms_accepted FROM users WHERE id = ?", (session["user_id"],)).fetchone()
        if not user or not user["terms_accepted"]:
            return redirect(url_for("terms_page"))
        return f(*args, **kwargs)
    return wrapper


def admin_required(f):
    @wraps(f)
    def wrapper(*args, **kwargs):
        key = request.args.get("key") or request.headers.get("X-Admin-Key")
        if not ADMIN_SECRET or key != ADMIN_SECRET:
            return "Access denied.", 403
        return f(*args, **kwargs)
    return wrapper


# =========================
# WELCOME / INTRODUCTION
# =========================



# ============================================================

# =========================
# NOVARA USAGE LIMITS
# =========================

NOVARA_USAGE_LIMITS = {
    "free": {
        "image": 1,
        "silent_video": 1,
        "audio_video": 0,
    },
    "plus": {
        "image": 40,
        "silent_video": 5,
        "audio_video": 5,
    },
    "pro": {
        "image": 80,
        "silent_video": 10,
        "audio_video": 10,
    },
    "v3.2": {
        "image": 120,
        "silent_video": 15,
        "audio_video": 15,
    },
}


def novara_usage_limits(plan):
    plan = str(plan or "free").strip().lower()
    return NOVARA_USAGE_LIMITS.get(plan, NOVARA_USAGE_LIMITS["free"])



def novara_next_upgrade_plan(plan):
    current = str(plan or "free").strip().lower()
    plan_order = [str(name).strip().lower() for name in NOVARA_USAGE_LIMITS.keys()]

    try:
        index = plan_order.index(current)
    except ValueError:
        index = 0

    if index + 1 >= len(plan_order):
        return None

    return plan_order[index + 1]


def novara_upgrade_limit_response(usage_type, plan, used, limit):
    current = str(plan or "free").strip().lower()
    next_plan = novara_next_upgrade_plan(current)

    labels = {
        "image": "image generation",
        "silent_video": "video generation",
        "audio_video": "audio + video generation",
    }

    feature = labels.get(usage_type, "this feature")

    if next_plan:
        message = (
            f"Your monthly {feature} limit has been reached ({used}/{limit}). "
            f"You're currently on the {current} plan. "
            f"Upgrade to the {next_plan} plan to get a higher monthly limit."
        )
    else:
        message = (
            f"Your monthly {feature} limit has been reached ({used}/{limit}). "
            f"You're already on the highest available plan ({current})."
        )

    return {
        "error": message,
        "upgrade_required": bool(next_plan),
        "current_plan": current,
        "next_plan": next_plan,
        "usage_type": usage_type,
        "used": used,
        "limit": limit,
    }

def reset_monthly_usage_if_needed(db, user_id, user):
    current_month = datetime.now().strftime("%Y-%m")

    if user["usage_month"] != current_month:
        db.execute(
            """
            UPDATE users
            SET image_used = 0,
                silent_video_used = 0,
                audio_video_used = 0,
                usage_month = ?
            WHERE id = ?
            """,
            (current_month, user_id),
        )
        db.commit()

        user = db.execute(
            """
            SELECT plan, image_used, silent_video_used,
                   audio_video_used, usage_month
            FROM users
            WHERE id = ?
            """,
            (user_id,),
        ).fetchone()

    return user


def check_novara_usage(db, user, usage_type):
    limits = novara_usage_limits(user["plan"])

    column_map = {
        "image": "image_used",
        "silent_video": "silent_video_used",
        "audio_video": "audio_video_used",
    }

    column = column_map.get(usage_type)

    if not column:
        return False, 0, 0

    used = int(user[column] or 0)
    limit = int(limits.get(usage_type, 0))

    return used < limit, used, limit


def increment_novara_usage(db, user_id, usage_type):
    column_map = {
        "image": "image_used",
        "silent_video": "silent_video_used",
        "audio_video": "audio_video_used",
    }

    column = column_map.get(usage_type)

    if not column:
        raise ValueError("Invalid Novara usage type")

    db.execute(
        f"UPDATE users SET {column} = COALESCE({column}, 0) + 1 WHERE id = ?",
        (user_id,),
    )
    db.commit()




# ============================================================
# NOVARA REWARDED AD MISSIONS
# ============================================================

# Ad rewards operate in rolling 6-hour windows.
REWARDED_MISSION_WINDOW_HOURS = 0  # 0 = no reset/cooldown
REWARDED_MISSION_AD_LIMIT = 0  # 0 = unlimited rewarded ads


def rewarded_window_reset_if_needed(db, user_id, user):
    """
    Reset the rewarded-ad mission window every 6 hours.

    Normal monthly image/video usage is intentionally NOT reset here.
    Only the rewarded mission counters/credits are managed separately.
    """
    now = datetime.now()

    raw = user["rewarded_window_start"]

    reset = False

    if not raw:
        reset = True
    else:
        try:
            start = datetime.fromisoformat(str(raw))
            elapsed = now - start
            if REWARDED_MISSION_WINDOW_HOURS > 0 and elapsed.total_seconds() >= REWARDED_MISSION_WINDOW_HOURS * 3600:
                reset = True
        except Exception:
            reset = True

    if reset:
        window = now.isoformat()

        db.execute(
            """
            UPDATE users
            SET rewarded_window_start = ?,
                rewarded_today = 0
            WHERE id = ?
            """,
            (window, user_id),
        )

        user = db.execute(
            """
            SELECT plan,
                   image_used,
                   silent_video_used,
                   audio_video_used,
                   rewarded_today,
                   rewarded_window_start,
                   rewarded_image_credits,
                   rewarded_audio_video_credits,
                   usage_month
            FROM users
            WHERE id = ?
            """,
            (user_id,),
        ).fetchone()

    return user



def ensure_rewarded_claim_columns(db):
    """Add idempotency fields to existing SQLite users tables if needed."""
    columns = {
        row[1]
        for row in db.execute("PRAGMA table_info(users)").fetchall()
    }

    additions = [
        ("rewarded_last_claim_id", "TEXT"),
        ("rewarded_last_claim_mission", "TEXT"),
        ("rewarded_last_image_delta", "INTEGER DEFAULT 0"),
        ("rewarded_last_audio_video_delta", "INTEGER DEFAULT 0"),
    ]

    for name, definition in additions:
        if name not in columns:
            db.execute(
                f"ALTER TABLE users ADD COLUMN {name} {definition}"
            )


def ensure_rewarded_claims_table(db):
    """
    Persist every rewarded-ad claim ID permanently.

    A primary key prevents the same successful ad callback
    from being rewarded twice, even after later claims occur.
    """
    db.execute("""
        CREATE TABLE IF NOT EXISTS rewarded_ad_claims (
            claim_id TEXT PRIMARY KEY,
            user_id INTEGER NOT NULL,
            mission TEXT NOT NULL,
            rewarded_ads INTEGER NOT NULL,
            image_credits_added INTEGER NOT NULL DEFAULT 0,
            audio_video_credits_added INTEGER NOT NULL DEFAULT 0,
            created_at TEXT NOT NULL
        )
    """)
    db.commit()


@app.route("/api/app-config", methods=["GET"])
def app_config():
    try:
        version_path = os.path.join(BASE_DIR, "version.json")

        with open(version_path, "r", encoding="utf-8") as f:
            v = json.load(f)

        return jsonify({
            "versionCode": int(v.get("versionCode", 1)),
            "versionName": str(v.get("versionName", "1.0")),
            "isMajor": bool(v.get("isMajor", False)),
            "apkUrl": str(v.get("apkUrl", "")),
            "ui": v.get("ui", {}),
            "features": v.get("features", {}),
            "changelog": v.get("changelog", [])
        })
    except Exception as e:
        return jsonify({
            "versionCode": 1,
            "versionName": "1.0",
            "isMajor": False,
            "apkUrl": "",
            "ui": {},
            "features": {},
            "changelog": [],
            "error": str(e)
        }), 200


@app.route("/api/rewarded-ad/claim", methods=["POST"])
@login_required
def api_rewarded_ad_claim():
    """
    Claim one completed rewarded-ad callback.

    The Android client calls this only after the rewarded-ad
    SDK reports that the user earned the reward.

    The server alone decides the reward.

    Idempotency:
      - every successful callback has a unique claim_id
      - every processed claim_id is stored permanently
      - retrying the same claim_id returns its original result
      - old claim IDs cannot become rewards again
    """
    db = get_db()
    ensure_rewarded_claim_columns(db)
    ensure_rewarded_claims_table(db)

    payload = request.get_json(silent=True) or {}
    claim_id = str(payload.get("claim_id") or "").strip()

    if not claim_id or len(claim_id) > 128:
        return jsonify({
            "ok": False,
            "error": "Invalid rewarded-ad claim."
        }), 400

    # Check permanent claim history first.
    existing = db.execute(
        """
        SELECT mission,
               rewarded_ads,
               image_credits_added,
               audio_video_credits_added
        FROM rewarded_ad_claims
        WHERE claim_id = ?
          AND user_id = ?
        """,
        (claim_id, session["user_id"]),
    ).fetchone()

    if existing:
        return jsonify({
            "ok": True,
            "idempotent": True,
            "mission": existing["mission"],
            "rewarded_ads": int(existing["rewarded_ads"]),
            "limit": REWARDED_MISSION_AD_LIMIT,
            "image_credits_added": int(
                existing["image_credits_added"]
            ),
            "audio_video_credits_added": int(
                existing["audio_video_credits_added"]
            ),
            "message": "Reward already processed."
        })

    user = db.execute(
        """
        SELECT plan,
               image_used,
               silent_video_used,
               audio_video_used,
               rewarded_today,
               rewarded_window_start,
               rewarded_image_credits,
               rewarded_audio_video_credits,
               usage_month
        FROM users
        WHERE id = ?
        """,
        (session["user_id"],),
    ).fetchone()

    if not user:
        return jsonify({
            "error": "User account not found."
        }), 401

    # Serialize the mission update so two simultaneous claims
    # cannot both receive the same mission position.
    db.execute("BEGIN IMMEDIATE")

    try:
        # Re-check permanent history after obtaining the write lock.
        existing = db.execute(
            """
            SELECT mission,
                   rewarded_ads,
                   image_credits_added,
                   audio_video_credits_added
            FROM rewarded_ad_claims
            WHERE claim_id = ?
              AND user_id = ?
            """,
            (claim_id, session["user_id"]),
        ).fetchone()

        if existing:
            db.rollback()
            return jsonify({
                "ok": True,
                "idempotent": True,
                "mission": existing["mission"],
                "rewarded_ads": int(existing["rewarded_ads"]),
                "limit": REWARDED_MISSION_AD_LIMIT,
                "image_credits_added": int(
                    existing["image_credits_added"]
                ),
                "audio_video_credits_added": int(
                    existing["audio_video_credits_added"]
                ),
                "message": "Reward already processed."
            })

        locked = db.execute(
            """
            SELECT rewarded_today
            FROM users
            WHERE id = ?
            """,
            (session["user_id"],),
        ).fetchone()

        completed = int(locked["rewarded_today"] or 0)

        if REWARDED_MISSION_AD_LIMIT > 0 and completed >= REWARDED_MISSION_AD_LIMIT:
            db.rollback()
            return jsonify({
                "ok": False,
                "error": "Rewarded-ad mission limit reached.",
                "rewarded_ads": completed,
                "limit": REWARDED_MISSION_AD_LIMIT,
            }), 429

        # SERVER DECIDES THE REWARD.
        next_count = completed + 1
        cycle_position = ((next_count - 1) % 5) + 1

        image_credit_delta = 0
        audio_video_credit_delta = 0
        mission = ""

        if cycle_position == 1:
            image_credit_delta = 1
            mission = "image_1"
        elif cycle_position == 2:
            image_credit_delta = 2
            mission = "image_2"
        elif cycle_position == 3:
            image_credit_delta = 5
            mission = "image_3"
        elif cycle_position == 4:
            image_credit_delta = 5
            mission = "image_4"
        elif cycle_position == 5:
            image_credit_delta = 5
            mission = "image_5"


        db.execute(
            """
            UPDATE users
            SET rewarded_today = COALESCE(rewarded_today, 0) + 1,
                rewarded_image_credits =
                    COALESCE(rewarded_image_credits, 0) + ?,
                rewarded_audio_video_credits =
                    COALESCE(rewarded_audio_video_credits, 0) + ?
            WHERE id = ?
            """,
            (
                image_credit_delta,
                audio_video_credit_delta,
                session["user_id"],
            ),
        )

        db.execute(
            """
            INSERT INTO rewarded_ad_claims (
                claim_id,
                user_id,
                mission,
                rewarded_ads,
                image_credits_added,
                audio_video_credits_added,
                created_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (
                claim_id,
                session["user_id"],
                mission,
                next_count,
                image_credit_delta,
                audio_video_credit_delta,
                datetime.now().isoformat(),
            ),
        )

        db.commit()

    except Exception:
        db.rollback()
        raise

    return jsonify({
        "ok": True,
        "idempotent": False,
        "mission": mission,
        "rewarded_ads": next_count,
        "limit": REWARDED_MISSION_AD_LIMIT,
        "image_credits_added": image_credit_delta,
        "audio_video_credits_added": audio_video_credit_delta,
        "message": "Reward earned successfully."
    })


@app.route("/api/rewarded-ad/status", methods=["GET"])
@login_required
def api_rewarded_ad_status():
    db = get_db()

    user = db.execute(
        """
        SELECT plan,
               rewarded_today,
               rewarded_window_start,
               rewarded_image_credits,
               rewarded_audio_video_credits
        FROM users
        WHERE id = ?
        """,
        (session["user_id"],),
    ).fetchone()

    if not user:
        return jsonify({"error": "User account not found."}), 401

    completed = int(user["rewarded_today"] or 0)

    return jsonify({
        "ok": True,
        "rewarded_ads": completed,
        "limit": REWARDED_MISSION_AD_LIMIT,
        "remaining": -1 if REWARDED_MISSION_AD_LIMIT <= 0 else max(
            0,
            REWARDED_MISSION_AD_LIMIT - completed
        ),
        "window_hours": REWARDED_MISSION_WINDOW_HOURS,
        "window_start": user["rewarded_window_start"],
        "image_credits": int(
            user["rewarded_image_credits"] or 0
        ),
        "audio_video_credits": int(
            user["rewarded_audio_video_credits"] or 0
        ),
    })



# =========================
# NOVARA USAGE API
# =========================


@app.route("/api/billing/verify-google", methods=["POST"])
def verify_google_billing():
    if not session.get("user_id"):
        return jsonify({"error": "Authentication required."}), 401

    data = request.get_json(silent=True) or {}

    purchase_token = str(
        data.get("purchase_token", "")
    ).strip()

    product_id = str(
        data.get("product_id", "")
    ).strip()

    package_name = str(
        data.get("package_name", "com.novara.app")
    ).strip()

    if not purchase_token or not product_id:
        return jsonify({
            "error": "purchase_token and product_id are required."
        }), 400

    if product_id not in NOVARA_GOOGLE_PRODUCTS:
        return jsonify({
            "error": "Unknown Google Play product."
        }), 400

    verification = verify_google_play_subscription(
        package_name,
        purchase_token,
        product_id,
    )

    if not verification.get("verified"):
        return jsonify({
            "ok": False,
            "verified": False,
            "error": verification.get(
                "error",
                "Google Play purchase could not be verified.",
            ),
        }), 402

    db = get_db()

    db.execute(
        """
        UPDATE users
        SET plan = ?,
            subscription_product_id = ?,
            subscription_expiry = ?,
            subscription_purchase_token = ?
        WHERE id = ?
        """,
        (
            verification["plan"],
            product_id,
            verification.get("expiry"),
            purchase_token,
            session["user_id"],
        ),
    )

    db.commit()

    return jsonify({
        "ok": True,
        "verified": True,
        "plan": verification["plan"],
        "product_id": product_id,
        "expiry": verification.get("expiry"),
        "order_id": verification.get("order_id"),
    }), 200

@app.route("/api/usage", methods=["GET"])
@login_required
def api_usage():
    db = get_db()

    user = db.execute(
        """
        SELECT plan, image_used, silent_video_used,
               audio_video_used, rewarded_today,
               rewarded_window_start,
               rewarded_image_credits,
               rewarded_audio_video_credits,
               usage_month
        FROM users
        WHERE id = ?
        """,
        (session["user_id"],),
    ).fetchone()

    if not user:
        return jsonify({"error": "User account not found."}), 401

    user = reset_monthly_usage_if_needed(
        db,
        session["user_id"],
        user,
    )

    limits = novara_usage_limits(user["plan"])

    image_used = int(user["image_used"] or 0)
    silent_used = int(user["silent_video_used"] or 0)
    audio_used = int(user["audio_video_used"] or 0)
    rewarded_used = int(user["rewarded_today"] or 0)

    image_limit = int(limits["image"])
    silent_limit = int(limits["silent_video"])
    audio_limit = int(limits["audio_video"])

    return jsonify({
        "plan": str(user["plan"] or "free"),
        "month": user["usage_month"],
        "images": {
            "used": image_used,
            "limit": image_limit,
            "remaining": max(0, image_limit - image_used),
        },
        "silent_videos": {
            "used": silent_used,
            "limit": silent_limit,
            "remaining": max(0, silent_limit - silent_used),
        },
        "audio_video": {
            "used": audio_used,
            "limit": audio_limit,
            "remaining": max(0, audio_limit - audio_used),
        },
        "rewarded_ads": {
            "used": rewarded_used,
            "limit": REWARDED_MISSION_AD_LIMIT,
            "remaining": -1 if REWARDED_MISSION_AD_LIMIT <= 0 else max(
                0,
                REWARDED_MISSION_AD_LIMIT - rewarded_used
            ),
            "window_hours": REWARDED_MISSION_WINDOW_HOURS,
            "window_start": user["rewarded_window_start"],
            "image_credits": int(
                user["rewarded_image_credits"] or 0
            ),
            "audio_video_credits": int(
                user["rewarded_audio_video_credits"] or 0
            ),
        },
    })


# NOVARA EMAIL VERIFICATION
# ============================================================

SMTP_HOST = os.getenv("SMTP_HOST", "smtp.gmail.com")
SMTP_PORT = int(os.getenv("SMTP_PORT", "587"))
SMTP_USERNAME = os.getenv("SMTP_USERNAME", "")
SMTP_PASSWORD = os.getenv("SMTP_PASSWORD", "")
MAIL_FROM = os.getenv("MAIL_FROM", SMTP_USERNAME)

VERIFICATION_CODE_MINUTES = 10


def send_verification_email(to_email, first_name, code):
    if not SMTP_USERNAME or not SMTP_PASSWORD:
        raise RuntimeError("SMTP email settings are not configured.")

    msg = EmailMessage()
    msg["Subject"] = "Your Novara verification code"
    msg["From"] = MAIL_FROM
    msg["To"] = to_email

    msg.set_content(
        f"""Hi {first_name},

Your Novara verification code is:

{code}

This code expires in 10 minutes.

If you did not create a Novara account, you can safely ignore this email.

— Novara
"""
    )

    with smtplib.SMTP(SMTP_HOST, SMTP_PORT, timeout=30) as smtp:
        smtp.starttls()
        smtp.login(SMTP_USERNAME, SMTP_PASSWORD)
        smtp.send_message(msg)


@app.route("/welcome")
def welcome_page():
    if session.get("user_id"):
        return redirect(url_for("index"))
    return render_template("welcome.html")


@app.route("/api/guest", methods=["POST"])
def api_guest():
    """Lets someone use Novara without creating an account. A lightweight
    guest user is created behind the scenes so conversations still save
    normally for this browser session."""
    db = get_db()
    guest_name = "Guest-" + uuid.uuid4().hex[:6]
    user_id = str(uuid.uuid4())
    db.execute(
        "INSERT INTO users (id, username, terms_accepted, created_at) VALUES (?, ?, 0, ?)",
        (user_id, guest_name, datetime.now().isoformat())
    )
    db.commit()
    session["user_id"] = user_id
    session["username"] = guest_name
    return jsonify({"ok": True, "redirect": url_for("terms_page")})


# =========================
# TERMS & CONDITIONS
# =========================

@app.route("/terms")
@login_required
def terms_page():
    db = get_db()
    user = db.execute("SELECT terms_accepted FROM users WHERE id = ?", (session["user_id"],)).fetchone()
    if user and user["terms_accepted"]:
        return redirect(url_for("index"))
    return render_template("terms.html")


@app.route("/api/accept-terms", methods=["POST"])
@login_required
def api_accept_terms():
    db = get_db()
    db.execute("UPDATE users SET terms_accepted = 1 WHERE id = ?", (session["user_id"],))
    db.commit()
    return jsonify({"ok": True})


# =========================
# AUTH ROUTES — password based
# =========================

@app.route("/signup", methods=["GET"])
def signup_page():
    if session.get("user_id"):
        return redirect(url_for("index"))
    return render_template("signup.html", google_enabled=bool(GOOGLE_CLIENT_ID))


@app.route("/api/signup", methods=["POST"])
def api_signup():
    data = request.get_json(silent=True) or {}

    first_name = (data.get("first_name") or "").strip()
    last_name = (data.get("last_name") or "").strip()
    username = (data.get("username") or "").strip()
    email = (data.get("email") or "").strip().lower()
    phone_number = (data.get("phone_number") or "").strip()
    password = data.get("password") or ""
    confirm_password = data.get("confirm_password") or ""

    # Required fields
    if not first_name:
        return jsonify({"error": "First name is required."}), 400

    if not last_name:
        return jsonify({"error": "Last name is required."}), 400

    if not username:
        return jsonify({"error": "Username is required."}), 400

    if not email:
        return jsonify({"error": "Email address is required."}), 400

    if not password:
        return jsonify({"error": "Password is required."}), 400

    if not confirm_password:
        return jsonify({"error": "Please confirm your password."}), 400

    # Name validation
    if len(first_name) > 50 or len(last_name) > 50:
        return jsonify({"error": "Name is too long."}), 400

    if not re.fullmatch(r"[A-Za-zÀ-ÿ' -]+", first_name):
        return jsonify({"error": "Please enter a valid first name."}), 400

    if not re.fullmatch(r"[A-Za-zÀ-ÿ' -]+", last_name):
        return jsonify({"error": "Please enter a valid last name."}), 400

    # Username
    if not re.fullmatch(r"[A-Za-z0-9_.-]{3,30}", username):
        return jsonify({
            "error": "Username must be 3–30 characters and use only letters, numbers, _ . or -."
        }), 400

    # Email
    email_pattern = (
        r"^[A-Za-z0-9.!#$%&'*+/=?^_`{|}~-]+@"
        r"[A-Za-z0-9](?:[A-Za-z0-9-]{0,61}[A-Za-z0-9])?"
        r"(?:\.[A-Za-z0-9](?:[A-Za-z0-9-]{0,61}[A-Za-z0-9])?)+$"
    )

    if not re.fullmatch(email_pattern, email):
        return jsonify({"error": "Please enter a valid email address."}), 400

    # Optional phone
    if phone_number:
        if not re.fullmatch(r"[0-9+\-() ]{7,20}", phone_number):
            return jsonify({"error": "Please enter a valid phone number."}), 400

    # Password confirmation
    if password != confirm_password:
        return jsonify({"error": "Passwords do not match."}), 400

    # Common strong password requirements
    if len(password) < 8:
        return jsonify({
            "error": "Password must be at least 8 characters."
        }), 400

    if not re.search(r"[A-Z]", password):
        return jsonify({
            "error": "Password must contain an uppercase letter."
        }), 400

    if not re.search(r"[a-z]", password):
        return jsonify({
            "error": "Password must contain a lowercase letter."
        }), 400

    if not re.search(r"[0-9]", password):
        return jsonify({
            "error": "Password must contain a number."
        }), 400

    if not re.search(r"[^A-Za-z0-9]", password):
        return jsonify({
            "error": "Password must contain a special character."
        }), 400

    # NOVARA_COMMON_PASSWORD_CHECK
    # Reject common/easily guessed passwords.
    common_passwords = {
        "password", "password1", "password123",
        "12345678", "123456789", "1234567890",
        "qwerty", "qwerty123", "abc123",
        "letmein", "welcome", "welcome123",
        "admin", "admin123", "iloveyou",
        "monkey", "dragon", "football",
        "00000000", "11111111"
    }

    password_lower = password.lower()

    if password_lower in common_passwords:
        return jsonify({
            "error": "This password is too common. Please choose a stronger password."
        }), 400

    if password_lower == username.lower():
        return jsonify({
            "error": "Password cannot be the same as your username."
        }), 400

    if email and password_lower == email.split("@")[0].lower():
        return jsonify({
            "error": "Password cannot be based on your email address."
        }), 400

    # Reject commonly used passwords even when they satisfy complexity rules.
    common_passwords = {
        "password123!",
        "password123",
        "password1!",
        "password!",
        "qwerty123!",
        "qwerty123",
        "qwerty!",
        "welcome123!",
        "welcome123",
        "letmein123!",
        "letmein123",
        "admin123!",
        "admin123",
        "abc12345!",
        "abc12345",
        "iloveyou123!",
        "iloveyou123",
        "changeme123!",
        "changeme123",
        "novara123!",
        "novara123",
    }

    if password.lower() in common_passwords:
        return jsonify({
            "error": "That password is too common. Please choose a different password."
        }), 400

    db = get_db()

    # Existing username
    if db.execute(
        "SELECT id FROM users WHERE LOWER(username)=LOWER(?)",
        (username,)
    ).fetchone():
        return jsonify({
            "error": "That username is already taken."
        }), 400

    # Existing email
    if db.execute(
        "SELECT id FROM users WHERE LOWER(email)=LOWER(?)",
        (email,)
    ).fetchone():
        return jsonify({
            "error": "An account with that email already exists."
        }), 400

    # Six digit verification code
    verification_code = f"{secrets.randbelow(1000000):06d}"

    verification_code_hash = generate_password_hash(
        verification_code
    )

    verification_expires_at = (
        datetime.now() +
        timedelta(minutes=VERIFICATION_CODE_MINUTES)
    ).isoformat()

    user_id = str(uuid.uuid4())

    try:
        db.execute(
            """
            INSERT INTO users (
                id,
                username,
                email,
                password_hash,
                first_name,
                last_name,
                phone_number,
                email_verified,
                verification_code_hash,
                verification_expires_at,
                terms_accepted,
                created_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, 0, ?, ?, 0, ?)
            """,
            (
                user_id,
                username,
                email,
                generate_password_hash(password),
                first_name,
                last_name,
                phone_number or None,
                verification_code_hash,
                verification_expires_at,
                datetime.now().isoformat()
            )
        )

        db.commit()

        send_verification_email(
            email,
            first_name,
            verification_code
        )

    except Exception as e:
        db.rollback()
        print("NOVARA SIGNUP ERROR:", e)

        return jsonify({
            "error": "We couldn't send the verification email. Check your email settings and try again."
        }), 500

    session["user_id"] = user_id
    session["username"] = username

    return jsonify({
        "ok": True,
        "redirect": url_for("verify_email_page")
    })



@app.route("/login", methods=["GET"])
def login_page():
    if session.get("user_id"):
        return redirect(url_for("index"))
    return render_template("login.html", google_enabled=bool(GOOGLE_CLIENT_ID))


@app.route("/api/login", methods=["POST"])
def api_login():
    data = request.json or {}
    username = data.get("username", "").strip()
    password = data.get("password", "")

    db = get_db()
    user = db.execute("SELECT * FROM users WHERE username = ?", (username,)).fetchone()

    if not user or not user["password_hash"] or not check_password_hash(user["password_hash"], password):
        return jsonify({"error": "Username ya password galat hai."}), 401

    session["user_id"] = user["id"]
    session["username"] = user["username"] or user["email"]
    next_url = url_for("index") if user["terms_accepted"] else url_for("terms_page")
    return jsonify({"ok": True, "redirect": next_url})


# =========================
# AUTH ROUTES — Google OAuth
# =========================

@app.route("/auth/google")
def auth_google():
    if not GOOGLE_CLIENT_ID:
        return "Google login is not configured yet.", 400
    redirect_uri = url_for("auth_google_callback", _external=True)
    return oauth.google.authorize_redirect(redirect_uri)


@app.route("/auth/google/callback")
def auth_google_callback():
    token = oauth.google.authorize_access_token()
    userinfo = token.get("userinfo") or oauth.google.parse_id_token(token)

    google_id = userinfo["sub"]
    email = userinfo.get("email", "")
    name = userinfo.get("name", email.split("@")[0] if email else "user")

    db = get_db()
    user = db.execute("SELECT * FROM users WHERE google_id = ?", (google_id,)).fetchone()

    if not user:
        user_id = str(uuid.uuid4())
        db.execute(
            "INSERT INTO users (id, username, email, google_id, terms_accepted, created_at) VALUES (?, ?, ?, ?, 0, ?)",
            (user_id, name, email, google_id, datetime.now().isoformat())
        )
        db.commit()
        session["user_id"] = user_id
        session["username"] = name
        return redirect(url_for("terms_page"))
    else:
        session["user_id"] = user["id"]
        session["username"] = user["username"] or user["email"]
        return redirect(url_for("index") if user["terms_accepted"] else url_for("terms_page"))


@app.route("/api/auth/google-native", methods=["POST"])
def api_auth_google_native():
    data = request.get_json(silent=True) or {}
    token = data.get("id_token") or ""

    if not token:
        return jsonify({"error": "Missing token."}), 400

    if not GOOGLE_CLIENT_ID:
        return jsonify({"error": "Google login is not configured yet."}), 400

    try:
        idinfo = google_id_token.verify_oauth2_token(
            token,
            google_auth_requests.Request(),
            GOOGLE_CLIENT_ID
        )
    except Exception as e:
        print("GOOGLE NATIVE AUTH ERROR:", e)
        return jsonify({"error": "Invalid Google token."}), 400

    google_id = idinfo["sub"]
    email = idinfo.get("email", "")
    name = idinfo.get("name", email.split("@")[0] if email else "user")

    db = get_db()
    user = db.execute(
        "SELECT * FROM users WHERE google_id = ?",
        (google_id,)
    ).fetchone()

    if not user:
        user_id = str(uuid.uuid4())
        db.execute(
            "INSERT INTO users (id, username, email, google_id, terms_accepted, created_at) VALUES (?, ?, ?, ?, 0, ?)",
            (user_id, name, email, google_id, datetime.now().isoformat())
        )
        db.commit()
        session["user_id"] = user_id
        session["username"] = name
        return jsonify({"ok": True})
    else:
        session["user_id"] = user["id"]
        session["username"] = user["username"] or user["email"]
        return jsonify({"ok": True})



@app.route("/api/account/delete", methods=["POST"])
@login_required
def api_account_delete():
    data = request.get_json(silent=True) or {}
    password = data.get("password") or ""

    db = get_db()

    user = db.execute(
        "SELECT id, password_hash, google_id FROM users WHERE id = ?",
        (session["user_id"],)
    ).fetchone()

    if not user:
        session.clear()
        return jsonify({"error": "Account not found."}), 404

    # Password accounts must confirm their password.
    # Google-only accounts can delete without a password.
    if user["password_hash"]:
        if not password or not check_password_hash(user["password_hash"], password):
            return jsonify({"error": "Incorrect password."}), 401

    user_id = user["id"]

    try:
        # Delete messages belonging to this user's conversations first.
        db.execute(
            """
            DELETE FROM messages
            WHERE conversation_id IN (
                SELECT id FROM conversations WHERE user_id = ?
            )
            """,
            (user_id,)
        )

        # Delete the user's conversations.
        db.execute(
            "DELETE FROM conversations WHERE user_id = ?",
            (user_id,)
        )

        # Delete feedback submitted by the user.
        db.execute(
            "DELETE FROM feedback_events WHERE user_id = ?",
            (user_id,)
        )

        # Delete per-user settings.
        db.execute(
            "DELETE FROM user_settings WHERE user_id = ?",
            (user_id,)
        )

        # Finally delete the account.
        db.execute(
            "DELETE FROM users WHERE id = ?",
            (user_id,)
        )

        db.commit()

    except Exception:
        db.rollback()
        return jsonify({
            "error": "Account deletion failed. Please try again."
        }), 500

    session.clear()

    return jsonify({
        "ok": True,
        "redirect": url_for("account_deleted")
    })


@app.route("/account/delete-request")
def account_delete_request():
    db = get_db()

    logged_in = bool(session.get("user_id"))
    has_password = False

    if logged_in:
        user = db.execute(
            "SELECT password_hash FROM users WHERE id = ?",
            (session["user_id"],)
        ).fetchone()
        has_password = bool(user and user["password_hash"])

    return render_template(
        "account-delete.html",
        logged_in=logged_in,
        has_password=has_password,
        contact_email=os.getenv("CONTACT_EMAIL", "support@novara.app")
    )


@app.route("/account/deleted")
def account_deleted():
    return render_template("account-deleted.html")



@app.route("/api/me", methods=["GET"])
def api_me():
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not logged in"}), 401
    db = get_db()
    user = db.execute(
        "SELECT username, email, phone_number, google_id, created_at FROM users WHERE id = ?",
        (session["user_id"],)
    ).fetchone()
    if not user:
        return jsonify({"ok": False, "error": "User not found"}), 404
    return jsonify({
        "ok": True,
        "username": user["username"],
        "email": user["email"],
        "phone_number": user["phone_number"],
        "is_guest": user["email"] is None and user["google_id"] is None,
        "signed_in_with_google": user["google_id"] is not None,
        "created_at": user["created_at"]
    })


@app.route("/api/memories", methods=["GET"])
def api_list_memories():
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not logged in"}), 401
    db = get_db()
    rows = db.execute(
        "SELECT id, content, source, created_at FROM memories WHERE user_id = ? ORDER BY created_at DESC",
        (session["user_id"],)
    ).fetchall()
    return jsonify({
        "ok": True,
        "memories": [
            {"id": r["id"], "content": r["content"], "source": r["source"], "created_at": r["created_at"]}
            for r in rows
        ]
    })


@app.route("/api/memories/<memory_id>", methods=["DELETE"])
def api_delete_memory(memory_id):
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not logged in"}), 401
    db = get_db()
    owned = db.execute(
        "SELECT id FROM memories WHERE id = ? AND user_id = ?",
        (memory_id, session["user_id"])
    ).fetchone()
    if not owned:
        return jsonify({"ok": False, "error": "Memory not found"}), 404
    db.execute("DELETE FROM memories WHERE id = ?", (memory_id,))
    db.commit()
    return jsonify({"ok": True})


@app.route("/api/memories", methods=["DELETE"])
def api_clear_memories():
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not logged in"}), 401
    db = get_db()
    db.execute("DELETE FROM memories WHERE user_id = ?", (session["user_id"],))
    db.commit()
    return jsonify({"ok": True})


@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login_page"))


# =========================
# ADMIN PANEL
# =========================

@app.route("/admin")
@admin_required
def admin_panel():
    db = get_db()
    users = db.execute(
        "SELECT id, username, email, created_at FROM users ORDER BY created_at DESC"
    ).fetchall()
    result = []
    for u in users:
        conv_count = db.execute(
            "SELECT COUNT(*) as c FROM conversations WHERE user_id = ?", (u["id"],)
        ).fetchone()["c"]
        result.append({
            "username": u["username"] or "-",
            "email": u["email"] or "-",
            "created_at": u["created_at"],
            "conversations": conv_count
        })
    return render_template("admin.html", users=result)


# =========================
# PDF KNOWLEDGE (in-memory, populated as PDFs are attached)
# =========================

pdf_chunks = []



def extract_text_attachment(path, filename, ext):
    """Extract readable text/data from supported office/text attachments."""
    try:
        if ext == "txt":
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        if ext == "json":
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                data = json.load(f)
            return json.dumps(data, ensure_ascii=False, indent=2)

        if ext == "csv":
            rows = []
            with open(path, "r", encoding="utf-8-sig", errors="replace",
                      newline="") as f:
                reader = csv.reader(f)
                for row in reader:
                    rows.append(" | ".join(str(x) for x in row))
            return "\n".join(rows)

        if ext == "docx":
            if Document is None:
                return "[DOCX parser is not installed.]"
            doc = Document(path)
            parts = []

            for p in doc.paragraphs:
                if p.text.strip():
                    parts.append(p.text.strip())

            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(cell.text.strip()
                                            for cell in row.cells))

            return "\n".join(parts)

        if ext == "xlsx":
            if load_workbook is None:
                return "[XLSX parser is not installed.]"

            wb = load_workbook(path, read_only=True, data_only=True)
            parts = []

            for ws in wb.worksheets:
                parts.append(f"=== Sheet: {ws.title} ===")
                for row in ws.iter_rows(values_only=True):
                    values = [
                        "" if value is None else str(value)
                        for value in row
                    ]
                    if any(v.strip() for v in values):
                        parts.append(" | ".join(values))

            wb.close()
            return "\n".join(parts)

        if ext == "pptx":
            if Presentation is None:
                return "[PPTX parser is not installed.]"

            prs = Presentation(path)
            parts = []

            for number, slide in enumerate(prs.slides, 1):
                parts.append(f"=== Slide {number} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())

            return "\n".join(parts)

    except Exception as e:
        return f"[Could not read {filename}: {e}]"

    return ""


def index_pdf(path, filename):
    try:
        reader = PdfReader(path)
        for page_number, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if not text.strip():
                continue
            words = text.split()
            for i in range(0, len(words), 250):
                text_chunk = " ".join(words[i:i + 250])
                if text_chunk.strip():
                    pdf_chunks.append({"file": filename, "page": page_number, "text": text_chunk})
    except Exception as e:
        print("PDF INDEX ERROR:", e)


def search_pdf(question):
    question_words = set(re.findall(r"[a-zA-Z0-9]+", question.lower()))
    results = []
    for chunk in pdf_chunks:
        chunk_words = set(re.findall(r"[a-zA-Z0-9]+", chunk["text"].lower()))
        score = len(question_words & chunk_words)
        if score > 0:
            results.append((score, chunk))
    results.sort(key=lambda x: x[0], reverse=True)
    return [x[1] for x in results[:5]]


# =========================
# WEB SEARCH (no API key needed — DuckDuckGo HTML results)
# For higher quality results, swap this out for a paid search API
# (Serper, Tavily, Bing) — see README "Web search" section.
# =========================

def search_web(query, max_results=5):
    try:
        resp = requests.get(
            f"https://html.duckduckgo.com/html/?q={quote(query)}",
            headers={"User-Agent": "Mozilla/5.0 (compatible; NovaraBot/1.0)"},
            timeout=8,
        )
        results = []
        # Lightweight scrape — no external HTML parser dependency required.
        blocks = re.findall(
            r'result__a"[^>]*>(.*?)</a>.*?result__snippet[^>]*>(.*?)</a>',
            resp.text, re.S
        )
        for title_html, snippet_html in blocks[:max_results]:
            title = re.sub("<[^<]+?>", "", title_html).strip()
            snippet = re.sub("<[^<]+?>", "", snippet_html).strip()
            if title:
                results.append({"title": title, "snippet": snippet})
        return results
    except Exception as e:
        print("WEB SEARCH ERROR:", e)
        return []


# =========================
# INTENT DETECTION — keyword based
# =========================

PDF_HINTS = ["pdf", "notes", "according to my", "in my book", "in the book", "mera pdf", "meri pdf", "syllabus", "chapter", "document"]
WEB_HINTS = ["latest", "today", "current", "news", "price", "score", "weather", "abhi", "aaj ka", "recent"]


def detect_intent(question):
    q = question.lower()
    if any(w in q for w in PDF_HINTS):
        return "PDF"
    if any(w in q for w in WEB_HINTS):
        return "WEB"
    return "NORMAL"


# =========================
# GEMINI CALL
# =========================

def call_deepseek(parts, timeout=60):
    if not DEEPSEEK_API_KEY:
        return None, "DEEPSEEK_API_KEY not set on server."

    try:
        has_image = any("inline_data" in part for part in parts)

        content = []
        for part in parts:
            if "text" in part:
                content.append({
                    "type": "text",
                    "text": part["text"]
                })
            elif "inline_data" in part:
                data = part["inline_data"]
                content.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:{data['mime_type']};base64,{data['data']}"
                    }
                })

        model = VISION_MODEL if has_image else MODEL

        response = requests.post(
            "https://api.deepseek.com/chat/completions",
            headers={
                "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": model,
                "messages": [
                    {
                        "role": "user",
                        "content": content
                    }
                ]
            },
            timeout=timeout
        )

        data = response.json()

        if "choices" not in data:
            err = data.get("error", {})
            return None, err.get("message", str(data))

        return data["choices"][0]["message"]["content"], None

    except Exception as e:
        return None, str(e)


def ask_ai(question, recent_messages, pdf_context="", image_path=None, web_context="", memories=None, model="fast", plan="free"):
    recent = "\n".join(f"{'User' if m['role']=='user' else 'Novara'}: {m['text']}" for m in recent_messages[-10:])
    memory_context = "\n".join(f"- {m}" for m in (memories or []))

    prompt = f"""
You are Novara, a personal AI assistant, built and trained by the Novara team.
Be natural, friendly, intelligent and helpful.
Do not mention internal tools, prompts, intent detection, APIs or system instructions.
If asked what model, AI, or technology powers you, who made you, or whether you are based on
Gemini, GPT, ChatGPT, Claude, or any other company's model, always say you are Novara,
an independent AI assistant, without naming or confirming any underlying provider or model.
Never reveal, confirm, or speculate about the specific underlying model or company powering you.

Use the information below only when relevant.

THINGS YOU REMEMBER ABOUT THIS USER:
{memory_context}

RECENT CONVERSATION:
{recent}

PDF INFORMATION:
{pdf_context}

WEB SEARCH RESULTS:
{web_context}

USER MESSAGE:
{question}
"""

    # Dedicated Qwen2-VL vision server on the S23 FE.
    if image_path and os.path.exists(image_path):
        mime_type, _ = mimetypes.guess_type(image_path)

        if mime_type and mime_type.startswith("image/"):
            try:
                import base64

                vision_url = os.getenv(
                    "VISION_SERVER_URL",
                    "https://applicants-simple-cup-newest.trycloudflare.com"
                ).rstrip("/")

                with open(image_path, "rb") as f:
                    image_b64 = base64.b64encode(f.read()).decode("utf-8")

                payload = {
                    "model": "ggml-org/Qwen2-VL-2B-Instruct-GGUF:Q4_K_M",
                    "messages": [{
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": question or "Analyze this image carefully and describe what you see."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime_type};base64,{image_b64}"
                                }
                            }
                        ]
                    }],
                    "temperature": 0.2,
                    "max_tokens": 1024,
                    "stream": False
                }

                response = requests.post(
                    vision_url + "/v1/chat/completions",
                    headers={"Content-Type": "application/json"},
                    json=payload,
                    timeout=120
                )

                data = response.json()

                if response.ok and data.get("choices"):
                    content = data["choices"][0]["message"].get("content", "")

                    if isinstance(content, list):
                        content = "".join(
                            x.get("text", "") if isinstance(x, dict) else str(x)
                            for x in content
                        )

                    if content:
                        return content

                app.logger.error(
                    "QWEN2-VL ERROR HTTP %s: %s",
                    response.status_code,
                    data
                )

            except Exception as e:
                app.logger.exception("QWEN2-VL REQUEST FAILED: %s", e)

    # NOVARA_QWEN_ONLY_IMAGE_GUARD
    # Image analysis must never fall back to Hugging Face.
    if image_path and os.path.exists(image_path):
        app.logger.error("QWEN2-VL failed; refusing HF image fallback.")
        return "AI vision error: Qwen2-VL vision server did not return a response."

    text, error = call_novara_model(
        prompt,
        mode=model,
        plan=plan,
    )

    if error:
        app.logger.error("NOVARA AI ERROR: %s", error)
        return f"AI backend error: {error}"

    return text


# =========================
# IMAGE GENERATION
# =========================

def generate_image(prompt):
    """Generate an image locally using SDXS via stable-diffusion.cpp."""
    prompt = (prompt or "").strip()
    if not prompt:
        return None, "Image prompt is empty."

    model_path = "/home/ubuntu/stable-diffusion.cpp/models/sdxs.safetensors"
    sd_cli = "/home/ubuntu/stable-diffusion.cpp/build/bin/sd-cli"

    if not os.path.isfile(model_path):
        return None, f"SDXS model not found: {model_path}"

    if not os.path.isfile(sd_cli):
        return None, f"sd-cli not found: {sd_cli}"

    enhanced_prompt = (
        "Create exactly the image described by the user. Preserve named people, "
        "places, objects, teams, clothing, and the requested setting. Do not "
        "replace the main subject with an unrelated person or object. Do not "
        "add unrelated characters. Follow the user's composition and scene closely. "
        "Generate a high-quality image. User request: " + prompt
    )

    try:
        import subprocess

        filename = f"{uuid.uuid4()}_generated.png"
        save_path = os.path.join(UPLOAD_FOLDER, filename)

        command = [
            sd_cli,
            "-m", model_path,
            "-p", enhanced_prompt,
            "-W", "512",
            "-H", "512",
            "-o", save_path
        ]

        result = subprocess.run(
            command,
            capture_output=True,
            text=True,
            timeout=120
        )

        if result.returncode != 0:
            error = (result.stderr or result.stdout or "Unknown SDXS error").strip()
            app.logger.error("SDXS ERROR: %s", error)
            return None, f"SDXS image generation failed: {error}"

        if not os.path.isfile(save_path):
            return None, "SDXS completed but did not create an image."

        return filename, None

    except subprocess.TimeoutExpired:
        return None, "SDXS image generation timed out."
    except Exception as e:
        app.logger.exception("SDXS image generation exception")
        return None, f"SDXS image generation failed: {e}"


# =========================
# MAIN APP ROUTES
# =========================

@app.route("/")
def root():
    if not session.get("user_id"):
        return redirect(url_for("welcome_page"))
    return redirect(url_for("index"))


@app.route("/app")
@login_required
@terms_required
def index():
    return render_template("index.html", username=session.get("username"))


@app.route("/uploads/<path:filename>")
@login_required
def serve_upload(filename):
    return send_from_directory(UPLOAD_FOLDER, filename)


@app.route("/api/conversations", methods=["GET"])
@login_required
def api_list_conversations():
    db = get_db()
    convs = db.execute(
        "SELECT id, title, created_at FROM conversations WHERE user_id = ? ORDER BY created_at DESC",
        (session["user_id"],)
    ).fetchall()
    return jsonify({"conversations": [dict(c) for c in convs]})


@app.route("/api/conversations/<conv_id>/messages", methods=["GET"])
@login_required
def api_get_messages(conv_id):
    db = get_db()
    conv = db.execute(
        "SELECT id FROM conversations WHERE id = ? AND user_id = ?", (conv_id, session["user_id"])
    ).fetchone()
    if not conv:
        return jsonify({"error": "Not found"}), 404
    rows = db.execute(
        "SELECT id, role, text, attachment_path, attachment_type, sources, feedback FROM messages WHERE conversation_id = ? ORDER BY created_at ASC",
        (conv_id,)
    ).fetchall()
    return jsonify({"messages": [dict(r) for r in rows]})


@app.route("/api/conversations/<conv_id>", methods=["DELETE"])
@login_required
def api_delete_conversation(conv_id):
    db = get_db()
    db.execute("DELETE FROM conversations WHERE id = ? AND user_id = ?", (conv_id, session["user_id"]))
    db.execute("DELETE FROM messages WHERE conversation_id = ?", (conv_id,))
    db.commit()
    return jsonify({"ok": True})


@app.route("/api/conversations/<conv_id>/rename", methods=["POST"])
@login_required
def api_rename_conversation(conv_id):
    data = request.json or {}
    title = (data.get("title") or "").strip()
    if not title:
        return jsonify({"error": "Naam khaali nahi ho sakta."}), 400
    title = title[:60]
    db = get_db()
    conv = db.execute(
        "SELECT id FROM conversations WHERE id = ? AND user_id = ?", (conv_id, session["user_id"])
    ).fetchone()
    if not conv:
        return jsonify({"error": "Not found"}), 404
    db.execute("UPDATE conversations SET title = ? WHERE id = ?", (title, conv_id))
    db.commit()
    return jsonify({"ok": True, "title": title})


@app.route("/api/conversations/<conv_id>/share", methods=["POST"])
@login_required
def api_share_conversation(conv_id):
    db = get_db()
    conv = db.execute(
        "SELECT id, share_id FROM conversations WHERE id = ? AND user_id = ?", (conv_id, session["user_id"])
    ).fetchone()
    if not conv:
        return jsonify({"error": "Not found"}), 404

    share_id = conv["share_id"] or uuid.uuid4().hex[:12]
    db.execute("UPDATE conversations SET share_id = ?, is_shared = 1 WHERE id = ?", (share_id, conv_id))
    db.commit()
    return jsonify({"ok": True, "share_url": url_for("view_shared", share_id=share_id, _external=True)})


@app.route("/api/conversations/<conv_id>/unshare", methods=["POST"])
@login_required
def api_unshare_conversation(conv_id):
    db = get_db()
    db.execute(
        "UPDATE conversations SET is_shared = 0 WHERE id = ? AND user_id = ?", (conv_id, session["user_id"])
    )
    db.commit()
    return jsonify({"ok": True})


@app.route("/share/<share_id>")
def view_shared(share_id):
    db = get_db()
    conv = db.execute(
        "SELECT id, title FROM conversations WHERE share_id = ? AND is_shared = 1", (share_id,)
    ).fetchone()
    if not conv:
        return render_template("share.html", not_found=True, title=None, messages=[])
    rows = db.execute(
        "SELECT role, text, attachment_type FROM messages WHERE conversation_id = ? ORDER BY created_at ASC",
        (conv["id"],)
    ).fetchall()
    return render_template("share.html", not_found=False, title=conv["title"], messages=[dict(r) for r in rows])


@app.route("/api/messages/<message_id>/feedback", methods=["POST"])
@login_required
def api_message_feedback(message_id):
    data = request.json or {}
    feedback = data.get("feedback")  # "like", "dislike", or null to clear
    if feedback not in ("like", "dislike", None):
        return jsonify({"error": "Invalid feedback"}), 400
    db = get_db()
    db.execute("UPDATE messages SET feedback = ? WHERE id = ?", (feedback, message_id))
    db.commit()
    return jsonify({"ok": True})



# NOVARA_FILE_READER_V2
def novara_read_uploaded_file(file_storage):
    import io
    import csv
    import json
    import zipfile
    import xml.etree.ElementTree as ET

    filename = str(getattr(file_storage, "filename", "") or "")
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    raw = file_storage.read()
    if not raw:
        return "", "The selected file is empty."

    # TXT / CSV / JSON
    if ext in {"txt", "csv", "json"}:
        text = None
        for enc in ("utf-8", "utf-8-sig", "utf-16", "latin-1"):
            try:
                text = raw.decode(enc)
                break
            except UnicodeDecodeError:
                pass

        if text is None:
            return "", "Could not decode the selected file."

        if ext == "json":
            try:
                text = json.dumps(
                    json.loads(text),
                    ensure_ascii=False,
                    indent=2
                )
            except Exception:
                pass

        elif ext == "csv":
            try:
                rows = csv.reader(io.StringIO(text))
                text = "\n".join(" | ".join(row) for row in rows)
            except Exception:
                pass

        return text[:120000], None

    # DOCX / XLSX / PPTX are ZIP/XML based formats.
    if ext not in {"docx", "xlsx", "pptx"}:
        return "", None

    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as z:

            # DOCX
            if ext == "docx":
                root = ET.fromstring(z.read("word/document.xml"))
                parts = []

                for e in root.iter():
                    if e.tag.endswith("}t") and e.text:
                        parts.append(e.text)
                    elif e.tag.endswith("}br"):
                        parts.append("\n")

                return " ".join(parts).strip()[:120000], None

            # XLSX
            if ext == "xlsx":
                shared = []

                if "xl/sharedStrings.xml" in z.namelist():
                    root = ET.fromstring(z.read("xl/sharedStrings.xml"))
                    for item in root:
                        parts = []
                        for e in item.iter():
                            if e.tag.endswith("}t") and e.text:
                                parts.append(e.text)
                        shared.append("".join(parts))

                sheets = []
                for name in sorted(z.namelist()):
                    if not name.startswith("xl/worksheets/sheet"):
                        continue
                    if not name.endswith(".xml"):
                        continue

                    root = ET.fromstring(z.read(name))
                    rows = []

                    for row in root.iter():
                        if not row.tag.endswith("}row"):
                            continue

                        values = []

                        for cell in row:
                            if not cell.tag.endswith("}c"):
                                continue

                            value = ""
                            cell_type = cell.attrib.get("t")

                            for child in cell:
                                if child.tag.endswith("}v") and child.text:
                                    value = child.text
                                elif child.tag.endswith("}is"):
                                    parts = []
                                    for e in child.iter():
                                        if e.tag.endswith("}t") and e.text:
                                            parts.append(e.text)
                                    value = "".join(parts)

                            if cell_type == "s":
                                try:
                                    value = shared[int(value)]
                                except Exception:
                                    pass

                            values.append(str(value))

                        if values:
                            rows.append(" | ".join(values))

                    if rows:
                        sheets.append("\n".join(rows))

                return "\n\n".join(sheets)[:120000], None

            # PPTX
            if ext == "pptx":
                slides = []

                for name in sorted(z.namelist()):
                    if not name.startswith("ppt/slides/slide"):
                        continue
                    if not name.endswith(".xml"):
                        continue

                    root = ET.fromstring(z.read(name))
                    parts = []

                    for e in root.iter():
                        if e.tag.endswith("}t") and e.text:
                            parts.append(e.text)

                    if parts:
                        slides.append(" ".join(parts))

                return "\n\n".join(slides)[:120000], None

    except Exception as exc:
        return "", f"Could not read {filename}: {exc}"

    return "", None

def generate_conversation_title(question, answer="", plan="free"):
    """Generate a short AI title for conversation history."""
    fallback = (question[:40] + "…") if len(question) > 40 else (question or "New chat")

    if not question:
        return "New chat"

    prompt = f"""Create a short title for this conversation.

Rules:
- 2 to 7 words
- Describe the main topic
- Natural for a chat history
- No quotes
- No punctuation at the end
- Return ONLY the title

USER:
{question[:1000]}

ASSISTANT:
{answer[:1000]}
"""

    try:
        title, error = call_novara_model(
            prompt,
            mode="fast",
            plan=plan,
            timeout=15,
        )

        if error or not title:
            return fallback

        title = " ".join(title.strip().split())
        title = title.strip('"').strip("'").rstrip(".").strip()

        return title[:80] if title else fallback
    except Exception:
        return fallback


@app.route("/api/chat", methods=["POST"])
@login_required
@terms_required
def api_chat():
    question = request.form.get("message", "").strip()
    conv_id = request.form.get("conversation_id") or None

    # Android sends: fast / thinking / omega
    requested_model = request.form.get("model", "fast").strip().lower()
    if requested_model not in {"fast", "thinking", "omega"}:
        requested_model = "fast"

    # NEVER trust the plan supplied by Android.
    # The database is the authority.
    db = get_db()
    current_user = db.execute(
        "SELECT plan FROM users WHERE id = ?",
        (session["user_id"],)
    ).fetchone()

    user_plan = (
        current_user["plan"]
        if current_user and current_user["plan"]
        else "free"
    )

    user_plan = str(user_plan).strip().lower()

    force_search = request.form.get("web_search") == "true"
    force_image_gen = request.form.get("image_gen") == "true"
    file = request.files.get("file")

    # NOVARA_FILE_READER_V2_ATTACH
    extracted_file_text = ""
    extracted_file_error = None

    if file and getattr(file, "filename", None):
        filename_lower = file.filename.lower()
        extension = (
            filename_lower.rsplit(".", 1)[-1]
            if "." in filename_lower else ""
        )

        if extension in {"docx", "xlsx", "pptx", "txt", "csv", "json"}:
            extracted_file_text, extracted_file_error = (
                novara_read_uploaded_file(file)
            )

            if extracted_file_error:
                return jsonify({"error": extracted_file_error}), 400

            if extracted_file_text.strip():
                question = (
                    question
                    + "\n\n[Attached file: "
                    + file.filename
                    + "]\n"
                    + extracted_file_text
                ).strip()

            # These formats have already been consumed as text.
            # Prevent later image/PDF handling from treating them
            # as binary media.
            file = None



    if not question and not file:
        return jsonify({"error": "Empty message"}), 400

    db = get_db()
    is_new_conversation = False

    if not conv_id:
        is_new_conversation = True
        conv_id = str(uuid.uuid4())
        title = "New chat"
        db.execute(
            "INSERT INTO conversations (id, user_id, title, created_at) VALUES (?, ?, ?, ?)",
            (conv_id, session["user_id"], title, datetime.now().isoformat())
        )
    else:
        conv = db.execute(
            "SELECT id FROM conversations WHERE id = ? AND user_id = ?", (conv_id, session["user_id"])
        ).fetchone()
        if not conv:
            return jsonify({"error": "Conversation not found"}), 404

    # ---- Image generation path ----
    if force_image_gen and question:
        # Server-side monthly image limit.
        # The user's plan is read from the database above;
        # Android cannot override it.
        db = get_db()

        usage_user = db.execute(
            """
            SELECT plan, image_used, silent_video_used,
                   audio_video_used, usage_month
            FROM users
            WHERE id = ?
            """,
            (session["user_id"],),
        ).fetchone()

        usage_user = reset_monthly_usage_if_needed(
            db,
            session["user_id"],
            usage_user,
        )

        allowed, used, limit = check_novara_usage(
            db,
            usage_user,
            "image",
        )

        if not allowed:
            return jsonify(
                novara_upgrade_limit_response(
                    "image",
                    usage_user["plan"],
                    used,
                    limit,
                )
            ), 429

        image_filename, gen_error = generate_image(question)
        now = datetime.now().isoformat()

        db.execute(
            "INSERT INTO messages (id, conversation_id, role, text, attachment_path, attachment_type, sources, created_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
            (str(uuid.uuid4()), conv_id, "user", question, None, None, "", now)
        )

        assistant_msg_id = str(uuid.uuid4())
        if gen_error:
            reply_text = f"I couldn't generate that image. ({gen_error})"
            db.execute(
                "INSERT INTO messages (id, conversation_id, role, text, attachment_path, attachment_type, sources, created_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                (assistant_msg_id, conv_id, "assistant", reply_text, None, None, "", now)
            )
            db.commit()
            return jsonify({"reply": reply_text, "message_id": assistant_msg_id, "conversation_id": conv_id})

        reply_text = "Here's what I created:"

        # Consume one image only after successful generation.
        increment_novara_usage(
            db,
            session["user_id"],
            "image",
        )

        db.execute(
            "INSERT INTO messages (id, conversation_id, role, text, attachment_path, attachment_type, sources, created_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
            (assistant_msg_id, conv_id, "assistant", reply_text, image_filename, "image", "", now)
        )
        db.commit()
        return jsonify({
            "reply": reply_text,
            "message_id": assistant_msg_id,
            "conversation_id": conv_id,
            "generated_image": image_filename
        })

    attachment_path = None
    attachment_type = None
    image_path_for_ai = None
    extracted_attachment_text = ""

    if file and file.filename:
        ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
        safe_name = f"{uuid.uuid4()}_{secure_filename(file.filename)}"
        save_path = os.path.join(UPLOAD_FOLDER, safe_name)
        file.save(save_path)
        attachment_path = safe_name

        if ext in IMAGE_EXT:
            attachment_type = "image"
            image_path_for_ai = save_path
        elif ext in VIDEO_EXT:
            attachment_type = "video"
        elif ext in DOC_EXT:
            attachment_type = "pdf"
            index_pdf(save_path, file.filename)
        elif ext in DOCX_EXT:
            attachment_type = "document"
            extracted_attachment_text = extract_text_attachment(
                save_path, file.filename, ext
            )
        elif ext in PPTX_EXT:
            attachment_type = "presentation"
            extracted_attachment_text = extract_text_attachment(
                save_path, file.filename, ext
            )
        elif ext in XLSX_EXT:
            attachment_type = "spreadsheet"
            extracted_attachment_text = extract_text_attachment(
                save_path, file.filename, ext
            )
        elif ext in TEXT_EXT:
            attachment_type = "text"
            extracted_attachment_text = extract_text_attachment(
                save_path, file.filename, ext
            )
        elif ext in PLUGIN_EXT:
            attachment_type = "file"
        else:
            attachment_type = "file"

    recent_rows = db.execute(
        "SELECT role, text FROM messages WHERE conversation_id = ? ORDER BY created_at DESC LIMIT 10",
        (conv_id,)
    ).fetchall()
    recent_messages = [dict(r) for r in reversed(recent_rows)]

    pdf_context = ""
    pdf_sources = []
    web_context = ""
    web_sources = []

    if question:
        intent = detect_intent(question)

        if intent == "PDF" or attachment_type == "pdf":
            results = search_pdf(question) if question else []
            if results:
                parts_txt = []
                for r in results:
                    parts_txt.append(f"Source: {r['file']}\nPage: {r['page']}\n\n{r['text']}\n")
                    pdf_sources.append({"file": r["file"], "page": r["page"]})
                pdf_context = "\n".join(parts_txt)

        if force_search or intent == "WEB":
            web_results = search_web(question)
            if web_results:
                web_context = "\n".join(f"- {r['title']}: {r['snippet']}" for r in web_results)
                web_sources = web_results

    effective_question = question if question else "(User attached a file — please respond to it.)"
    user_memories = get_user_memories(session["user_id"])
    if extracted_attachment_text:

        pdf_context = (

            (pdf_context + "\n\n" if pdf_context else "")

            + "ATTACHED FILE CONTENT:\n"

            + extracted_attachment_text[:120000]

        )

    try:
        answer = ask_ai(
            effective_question,
            recent_messages,
            pdf_context,
            image_path_for_ai,
            web_context,
            user_memories,
            model=requested_model,
            plan=user_plan,
        )
    except Exception as e:
        app.logger.exception("API CHAT AI ERROR")
        return jsonify({"error": f"AI backend error: {type(e).__name__}: {e}"}), 500

    if question:
        extract_and_save_memory(session["user_id"], question, answer, plan=user_plan)

    now = datetime.now().isoformat()
    db.execute(
        "INSERT INTO messages (id, conversation_id, role, text, attachment_path, attachment_type, sources, created_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
        (str(uuid.uuid4()), conv_id, "user", question, attachment_path, attachment_type, "", now)
    )
    assistant_msg_id = str(uuid.uuid4())
    combined_sources = ", ".join(f"{s['file']} p.{s['page']}" for s in pdf_sources)
    if web_sources:
        combined_sources = (combined_sources + " | " if combined_sources else "") + ", ".join(s["title"] for s in web_sources)

    db.execute(
        "INSERT INTO messages (id, conversation_id, role, text, attachment_path, attachment_type, sources, created_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
        (assistant_msg_id, conv_id, "assistant", answer, None, None, combined_sources, now)
    )
    db.commit()

    conversation_title = None

    if is_new_conversation:
        conversation_title = generate_conversation_title(
            question,
            answer,
            user_plan,
        )

        db.execute(
            "UPDATE conversations SET title = ? WHERE id = ? AND user_id = ?",
            (conversation_title, conv_id, session["user_id"]),
        )
        db.commit()

    return jsonify({
        "reply": answer,
        "conversation_title": conversation_title,
        "message_id": assistant_msg_id,
        "sources": pdf_sources,
        "web_sources": web_sources,
        "conversation_id": conv_id,
        "attachment_path": attachment_path,
        "attachment_type": attachment_type
    })


# =========================
# VIDEO STUDIO (generation + editing)
# Not wired to a provider yet — plug in VIDEO_API_URL / VIDEO_API_KEY
# once you pick a video-gen service. Until then these return a clear
# "not configured" message instead of failing silently.
# =========================

@app.route("/video")
@login_required
@terms_required
def video_studio():
    return render_template("video.html", username=session.get("username"), video_enabled=bool(VIDEO_API_URL))


@app.route("/api/video/generate", methods=["POST"])
@login_required
def api_video_generate():
    data = request.json or {}
    prompt = (data.get("prompt") or "").strip()
    video_type = str(data.get("video_type", "silent")).strip().lower()

    if video_type not in {"silent", "audio_video"}:
        video_type = "silent"

    if not prompt:
        return jsonify({"error": "Prompt cannot be empty."}), 400

    db = get_db()

    usage_user = db.execute(
        """
        SELECT plan, image_used, silent_video_used,
               audio_video_used, usage_month
        FROM users
        WHERE id = ?
        """,
        (session["user_id"],),
    ).fetchone()

    if not usage_user:
        return jsonify({"error": "User account not found."}), 401

    usage_user = reset_monthly_usage_if_needed(
        db,
        session["user_id"],
        usage_user,
    )

    usage_type = (
        "audio_video"
        if video_type == "audio_video"
        else "silent_video"
    )

    allowed, used, limit = check_novara_usage(
        db,
        usage_user,
        usage_type,
    )

    if not allowed:
        return jsonify(
            novara_upgrade_limit_response(
                usage_type,
                usage_user["plan"],
                used,
                limit,
            )
        ), 429

    replicate_token = os.getenv("REPLICATE_API_TOKEN")

    if not replicate_token:
        return jsonify({
            "error": "REPLICATE_API_TOKEN not set on server."
        }), 500

    # Resolution is ALWAYS decided server-side.
    # Free / Plus / Pro = 480p
    # V3.2 = 720p
    plan = str(usage_user["plan"] or "free").strip().lower()
    resolution = "720p" if plan in {"v3.2", "v3_2", "v32"} else "480p"

    model = "wan-video/wan-2.2-t2v-fast"

    try:
        import requests
        import time

        headers = {
            "Authorization": f"Bearer {replicate_token}",
            "Content-Type": "application/json",
        }

        payload = {
            "input": {
                "prompt": prompt,
                "go_fast": True,
                "num_frames": 81,
                "resolution": resolution,
                "aspect_ratio": "16:9",
            }
        }

        create_response = requests.post(
            f"https://api.replicate.com/v1/models/{model}/predictions",
            headers=headers,
            json=payload,
            timeout=60,
        )

        if create_response.status_code >= 400:
            return jsonify({
                "error": (
                    f"Replicate request failed "
                    f"({create_response.status_code}): "
                    f"{create_response.text[:1000]}"
                )
            }), 502

        prediction = create_response.json()
        prediction_url = prediction.get("urls", {}).get("get")

        if not prediction_url:
            return jsonify({
                "error": "Replicate did not return a prediction URL."
            }), 502

        # Wait for the generated MP4.
        for _ in range(120):
            time.sleep(5)

            status_response = requests.get(
                prediction_url,
                headers={
                    "Authorization": f"Bearer {replicate_token}"
                },
                timeout=30,
            )

            if status_response.status_code >= 400:
                return jsonify({
                    "error": (
                        f"Replicate status check failed "
                        f"({status_response.status_code})"
                    )
                }), 502

            result = status_response.json()
            status = result.get("status")

            if status == "succeeded":
                output = result.get("output")

                if isinstance(output, list):
                    output = output[0] if output else None

                if not output:
                    return jsonify({
                        "error": "Replicate completed without a video URL."
                    }), 502

                video_response = requests.get(
                    output,
                    timeout=300,
                )
                video_response.raise_for_status()

                filename = f"{uuid.uuid4()}_generated.mp4"
                save_path = os.path.join(
                    UPLOAD_FOLDER,
                    filename,
                )

                with open(save_path, "wb") as file:
                    file.write(video_response.content)

                # Consume allowance ONLY after successful generation.
                increment_novara_usage(
                    db,
                    session["user_id"],
                    usage_type,
                )

                return jsonify({
                    "ok": True,
                    "video": filename,
                    "video_url": url_for(
                        "serve_upload",
                        filename=filename,
                    ),
                    "model": model,
                    "resolution": resolution,
                    "video_type": video_type,
                })

            if status in {"failed", "canceled"}:
                error_detail = result.get("error") or "Unknown generation error."
                return jsonify({
                    "error": f"Video generation failed: {error_detail}"
                }), 502

        return jsonify({
            "error": "Video generation timed out."
        }), 504

    except requests.RequestException as e:
        return jsonify({
            "error": f"Replicate network error: {e}"
        }), 502

    except Exception as e:
        return jsonify({
            "error": f"Video generation failed: {e}"
        }), 500


@app.route("/api/video/edit", methods=["POST"])
@login_required
def api_video_edit():
    if not VIDEO_API_URL:
        return jsonify({"error": "Video editing is not configured yet. Set the VIDEO_API_URL and VIDEO_API_KEY environment variables. See the README for details."}), 501
    file = request.files.get("file")
    action = request.form.get("action", "enhance")
    if not file:
        return jsonify({"error": "No video file was provided."}), 400
    return jsonify({"error": "Video editing must be implemented for the selected provider. See the Video editing section of the README."}), 501


if __name__ == "__main__":
    port = int(os.getenv("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=False)



@app.route("/verify-email")
@login_required
def verify_email_page():
    db = get_db()

    user = db.execute(
        """
        SELECT email, first_name, email_verified
        FROM users
        WHERE id = ?
        """,
        (session["user_id"],)
    ).fetchone()

    if not user:
        session.clear()
        return redirect(url_for("signup_page"))

    if user["email_verified"]:
        return redirect(url_for("terms_page"))

    return render_template(
        "verify_email.html",
        email=user["email"],
        first_name=user["first_name"]
    )


@app.route("/api/verify-email", methods=["POST"])
@login_required
def api_verify_email():
    data = request.get_json(silent=True) or {}
    code = (data.get("code") or "").strip()

    if not re.fullmatch(r"\d{6}", code):
        return jsonify({
            "error": "Enter the 6-digit verification code."
        }), 400

    db = get_db()

    user = db.execute(
        """
        SELECT
            id,
            email_verified,
            verification_code_hash,
            verification_expires_at
        FROM users
        WHERE id = ?
        """,
        (session["user_id"],)
    ).fetchone()

    if not user:
        return jsonify({"error": "Account not found."}), 404

    if user["email_verified"]:
        return jsonify({
            "ok": True,
            "redirect": url_for("terms_page")
        })

    if not user["verification_code_hash"]:
        return jsonify({
            "error": "No verification code is available."
        }), 400

    try:
        expires_at = datetime.fromisoformat(
            user["verification_expires_at"]
        )

        if datetime.now() > expires_at:
            return jsonify({
                "error": "This verification code has expired."
            }), 400

    except Exception:
        return jsonify({
            "error": "Verification code is invalid."
        }), 400

    if not check_password_hash(
        user["verification_code_hash"],
        code
    ):
        return jsonify({
            "error": "Incorrect verification code."
        }), 400

    db.execute(
        """
        UPDATE users
        SET
            email_verified = 1,
            verification_code_hash = NULL,
            verification_expires_at = NULL
        WHERE id = ?
        """,
        (session["user_id"],)
    )

    db.commit()

    return jsonify({
        "ok": True,
        "redirect": url_for("terms_page")
    })


@app.route("/api/resend-verification", methods=["POST"])
@login_required
def api_resend_verification():
    db = get_db()

    user = db.execute(
        """
        SELECT id, email, first_name, email_verified
        FROM users
        WHERE id = ?
        """,
        (session["user_id"],)
    ).fetchone()

    if not user:
        return jsonify({"error": "Account not found."}), 404

    if user["email_verified"]:
        return jsonify({
            "ok": True,
            "redirect": url_for("terms_page")
        })

    code = f"{secrets.randbelow(1000000):06d}"

    code_hash = generate_password_hash(code)

    expires_at = (
        datetime.now() +
        timedelta(minutes=VERIFICATION_CODE_MINUTES)
    ).isoformat()

    try:
        send_verification_email(
            user["email"],
            user["first_name"],
            code
        )

        db.execute(
            """
            UPDATE users
            SET
                verification_code_hash = ?,
                verification_expires_at = ?
            WHERE id = ?
            """,
            (
                code_hash,
                expires_at,
                user["id"]
            )
        )

        db.commit()

        return jsonify({
            "ok": True,
            "message": "A new verification code has been sent."
        })

    except Exception as e:
        db.rollback()
        print("NOVARA RESEND ERROR:", e)

        return jsonify({
            "error": "Unable to send the verification email."
        }), 500



# =========================
# APP UPDATE
# =========================
@app.route("/download/latest")
def download_latest():
    return send_from_directory(BASE_DIR, "latest.apk", as_attachment=True)

@app.route("/api/app-version", methods=["GET"])
def app_version():
    try:
        version_path = os.path.join(BASE_DIR, "version.json")

        with open(
            version_path,
            "r",
            encoding="utf-8"
        ) as f:
            v = json.load(f)

        return jsonify({
            "versionCode": int(
                v.get("versionCode", 1)
            ),
            "versionName": str(
                v.get("versionName", "1.0")
            ),
            "isMajor": bool(
                v.get("isMajor", False)
            ),
            "apkUrl": str(
                v.get(
                    "apkUrl",
                    "/download/latest"
                )
            ),
            "changelog": v.get(
                "changelog",
                []
            )
        })

    except Exception as e:
        return jsonify({
            "versionCode": 1,
            "versionName": "1.0",
            "isMajor": False,
            "apkUrl": "/download/latest",
            "changelog": [],
            "error": str(e)
        }), 200
